"""Numeric stand-in for visual QA: LibreOffice here has no Impress filters,
so slides cannot be rendered. Liberation Sans is metric-compatible with Arial
(this template's body font), so text widths can be measured exactly instead.

Reports, per shape: estimated wrapped line count vs. available height, plus
slide-bounds and pairwise-overlap violations.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont

EMU_IN = 914400

#: Metric source for the template's body font. Liberation Sans is
#: metric-compatible with Arial; macOS ships Arial itself, which is better
#: still. Resolved at import so the checker runs on either machine rather than
#: only on the Linux box it was first written on -- an unrunnable layout check
#: is the same as no layout check.
_FONT_SETS = [
    {(False, False): "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
     (True,  False): "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
     (False, True):  "/usr/share/fonts/truetype/liberation/LiberationSans-Italic.ttf",
     (True,  True):  "/usr/share/fonts/truetype/liberation/LiberationSans-BoldItalic.ttf"},
    {(False, False): "/System/Library/Fonts/Supplemental/Arial.ttf",
     (True,  False): "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
     (False, True):  "/System/Library/Fonts/Supplemental/Arial Italic.ttf",
     (True,  True):  "/System/Library/Fonts/Supplemental/Arial Bold Italic.ttf"},
]
FONTS = next((f for f in _FONT_SETS if all(os.path.exists(p) for p in f.values())), None)
if FONTS is None:
    sys.exit("fitcheck: no Arial-metric font set found (Liberation Sans or macOS Arial)")
PX_PER_PT = 4.0   # render at 4x for sub-pixel accuracy, divide back out
_cache = {}

def font(sz_pt, bold, italic):
    key = (round(sz_pt, 1), bool(bold), bool(italic))
    if key not in _cache:
        _cache[key] = ImageFont.truetype(FONTS[(bool(bold), bool(italic))],
                                         max(1, int(round(sz_pt * PX_PER_PT))))
    return _cache[key]

def text_w_pt(s, sz, bold, italic):
    if not s: return 0.0
    return font(sz, bold, italic).getlength(s) / PX_PER_PT

def wrapped_lines(words_runs, avail_pt):
    """words_runs: list of (word, sz, bold, italic). Greedy wrap; returns line count
    and the max line width, plus the dominant size per line for height math."""
    lines, cur, curw, maxw = [], [], 0.0, 0.0
    for w, sz, b, i in words_runs:
        ww = text_w_pt(w + " ", sz, b, i)
        if cur and curw + ww > avail_pt:
            lines.append((cur, curw)); maxw = max(maxw, curw); cur, curw = [], 0.0
        cur.append((w, sz)); curw += ww
    if cur: lines.append((cur, curw)); maxw = max(maxw, curw)
    return lines, maxw

def check(path, only=None):
    prs = Presentation(path)
    SW, SH = prs.slide_width / EMU_IN, prs.slide_height / EMU_IN
    problems = []
    for si, slide in enumerate(prs.slides, 1):
        if only and si not in only: continue
        boxes = []
        for sh in slide.shapes:
            if sh.left is None: continue
            L, T = sh.left / EMU_IN, sh.top / EMU_IN
            W, H = (sh.width or 0) / EMU_IN, (sh.height or 0) / EMU_IN
            boxes.append((sh.name, L, T, W, H, sh))
            # slide bounds
            if L < -0.01 or T < -0.01 or L + W > SW + 0.01 or T + H > SH + 0.01:
                problems.append(f"s{si} {sh.name}: outside slide bounds "
                                f"({L:.2f},{T:.2f},{W:.2f}x{H:.2f})")
            if not sh.has_text_frame: continue
            tf = sh.text_frame
            if not (tf.text or "").strip():
                continue          # decorative shape: nothing to overflow
            li = (tf.margin_left or 0) / EMU_IN; ri = (tf.margin_right or 0) / EMU_IN
            ti = (tf.margin_top or 0) / EMU_IN;  bi = (tf.margin_bottom or 0) / EMU_IN
            avail_w_pt = max(1.0, (W - li - ri) * 72.0)
            total_h_pt = 0.0
            for para in tf.paragraphs:
                runs = []
                for r in para.runs:
                    sz = r.font.size.pt if r.font.size else 11.0
                    for w in (r.text or "").split():
                        runs.append((w, sz, r.font.bold, r.font.italic))
                if not runs:
                    continue      # blank paragraph: spacing handled by space_after
                lines, maxw = wrapped_lines(runs, avail_w_pt)
                for ln, _ in lines:
                    total_h_pt += max(s for _, s in ln) * 1.21
                total_h_pt += (para.space_after.pt if para.space_after else 0)
            avail_h_pt = max(1.0, (H - ti - bi) * 72.0)
            if total_h_pt > avail_h_pt * 1.02:
                over = total_h_pt - avail_h_pt
                problems.append(f"s{si} {sh.name}: TEXT OVERFLOW ~{over:.0f}pt "
                                f"(needs {total_h_pt:.0f}pt, box {avail_h_pt:.0f}pt)")
        # pairwise overlap among filled/text shapes
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                n1, l1, t1, w1, h1, s1 = boxes[a]; n2, l2, t2, w2, h2, s2 = boxes[b]
                if w1 <= 0 or h1 <= 0 or w2 <= 0 or h2 <= 0: continue
                ox = min(l1 + w1, l2 + w2) - max(l1, l2)
                oy = min(t1 + h1, t2 + h2) - max(t1, t2)
                if ox > 0.04 and oy > 0.04:
                    a1, a2 = w1 * h1, w2 * h2
                    contained = (ox * oy) >= 0.92 * min(a1, a2)
                    if contained:
                        continue   # nested label/card, not a collision
                    if ox * oy > 0.28 * min(a1, a2):
                        problems.append(f"s{si} OVERLAP {n1} x {n2} "
                                        f"({ox:.2f}x{oy:.2f}in)")
    return problems

if __name__ == "__main__":
    only = {int(x) for x in sys.argv[2:]} or None
    probs = check(sys.argv[1], only)
    if not probs:
        print("FIT CHECK: clean")
    else:
        print(f"FIT CHECK: {len(probs)} issue(s)")
        for p in probs: print("  " + p)
